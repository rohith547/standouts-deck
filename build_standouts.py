#!/usr/bin/env python3
"""Standout-projects speaker deck.

A focused deck of the flagship projects, each slide carrying a full ORATOR
SCRIPT in the speaker notes so the presenter can explain it like an expert:
hook -> story -> the hard part -> a landing line, plus likely Q&A.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Brand (matches the portfolio deck / Skills Hub) ------------------------
INK     = RGBColor(0x0B, 0x0E, 0x14)
PANEL   = RGBColor(0x11, 0x15, 0x1F)
EDGE    = RGBColor(0x1E, 0x25, 0x33)
ACCENT  = RGBColor(0x6E, 0xE7, 0xB7)   # emerald
ACCENT2 = RGBColor(0x38, 0xBD, 0xF8)   # sky
WHITE   = RGBColor(0xF3, 0xF4, 0xF6)
GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
DIM     = RGBColor(0x6B, 0x72, 0x80)

FONT = "Calibri"
NAME = "Rohith Shabad"
TAGLINE = "AI Application Engineer"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = INK
    bg.line.fill.background(); bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=True):
    t = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(t, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb


def chip_row(s, x0, y0, items, max_w, color=WHITE, gap=0.12, line_h=0.44):
    cx, cy = x0, y0
    for it in items:
        w = 0.20 + 0.092 * len(it)
        if cx + w > x0 + max_w:
            cx = x0; cy += line_h
        c = box(s, cx, cy, w, 0.34, fill=PANEL, line=EDGE, line_w=0.75)
        tf = c.text_frame; tf.word_wrap = False
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = it
        r.font.size = Pt(10.5); r.font.color.rgb = color; r.font.name = FONT
        cx += w + gap
    return cy + line_h


def header(s, kicker, title):
    box(s, 0.6, 0.62, 0.09, 0.62, fill=ACCENT)
    text(s, 0.85, 0.5, 11.8, 1.0, [
        [(kicker.upper(), 12, ACCENT, True)],
        [(title, 30, WHITE, True)],
    ], space_after=2)


def footer(s, n):
    text(s, 0.6, 7.02, 8, 0.35, [[(f"{NAME}  \u2022  {TAGLINE}", 9, DIM, False)]])
    text(s, 11.6, 7.02, 1.1, 0.35, [[(f"{n:02d}", 9, DIM, False)]], align=PP_ALIGN.RIGHT)


def notes(s, script):
    s.notes_slide.notes_text_frame.text = script


def deep_dive(n, kicker, title, repo, hook, what, how, hard, tags, script):
    s = slide(); header(s, kicker, title)
    # hook banner
    box(s, 0.85, 1.78, 11.6, 0.92, fill=PANEL, line=ACCENT, line_w=1.25)
    text(s, 1.15, 1.9, 11.0, 0.72, [[("\u201C" + hook + "\u201D", 16, WHITE, True)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.04)
    # three cards
    cards = [("What it does", what, ACCENT),
             ("How it works", how, ACCENT2),
             ("Why it's hard", hard, ACCENT)]
    xs = [0.85, 4.92, 8.99]
    for (t, d, col), x in zip(cards, xs):
        box(s, x, 2.95, 3.78, 2.55, fill=PANEL, line=EDGE)
        text(s, x + 0.28, 3.12, 3.25, 2.25, [
            [(t, 14, col, True)],
            [(d, 12, GRAY, False)],
        ], space_after=6, line_spacing=1.08)
    # tech + link
    text(s, 0.85, 5.72, 11.0, 0.35, [[
        ("Code  \u2192  ", 12, ACCENT2, True),
        (repo, 12, GRAY, False),
    ]])
    chip_row(s, 0.85, 6.12, tags, max_w=11.6, color=WHITE)
    footer(s, n)
    notes(s, script)


# ============================================================ 1. TITLE
s = slide()
box(s, 0, 0, 0.18, SH.inches, fill=ACCENT)
text(s, 1.0, 2.0, 11.5, 0.5, [[("STANDOUT PROJECTS  \u2022  SPEAKER DECK", 14, ACCENT, True)]])
text(s, 0.95, 2.45, 11.6, 2.0, [
    [(NAME, 50, WHITE, True)],
    [("The five projects I lead with \u2014 and how to talk about them.", 20, GRAY, False)],
], space_after=8)
text(s, 1.0, 4.7, 11.5, 0.6, [[
    ("Every slide has a full script in the ", 15, GRAY, False),
    ("Notes", 15, ACCENT, True),
    (" pane \u2014 hook, story, the hard part, and a line to land it.", 15, GRAY, False),
]])
notes(s, (
    "DELIVERY COACHING (read me once, then forget me):\n"
    "\u2022 Open with energy. Smile, slow down, and own the first sentence \u2014 the first 10 seconds set the tone.\n"
    "\u2022 Rhythm for each project: HOOK (one line) \u2192 pause \u2192 what it does \u2192 how it works \u2192 the HARD part \u2192 land it.\n"
    "\u2022 The 'hard part' is where you sound like an expert. Spend your energy there.\n"
    "\u2022 Pause after every bold claim. Silence makes you sound confident, not rushed.\n"
    "\u2022 Don't read bullets aloud \u2014 the slide is for them, the story is from you.\n"
    "\u2022 Close every project on the bold 'landing line' and stop talking. Let it sit.\n\n"
    "YOUR 30-SECOND ELEVATOR PITCH (memorize this):\n"
    "\u201CI'm an AI Application Engineer \u2014 I use AI to turn ideas into production apps fast. I've shipped 75-plus "
    "of them across five domains. Today I'll show you five standouts: an agent control center, an autonomous "
    "video pipeline, my own reusable-skills system, a suite of financial-intelligence dashboards, and a bank-"
    "automation bot that runs in production. The thread through all of them is the same \u2014 real, working systems, "
    "not demos.\u201D"
))

# ============================================================ 2. THE LINEUP
s = slide(); header(s, "Agenda", "The five standouts")
rows = [
    ("01  CommandCore", "Agent orchestration", "One chat that routes intent to specialized AI agents."),
    ("02  Beacon", "Content automation", "An AI news-video pipeline that ships itself daily."),
    ("03  Skills Hub", "Force-multiplier", "25 reusable skills wired into my AI command line."),
    ("04  Financial Intelligence", "Live data products", "Dashboards that turn market filings into signals."),
    ("05  Rent-tracker", "Production automation", "A bank-scraping bot that tracks who paid rent."),
]
y = 1.95
for tag, cat, d in rows:
    box(s, 0.85, y, 11.6, 0.86, fill=PANEL, line=EDGE)
    text(s, 1.15, y + 0.12, 3.3, 0.6, [[(tag, 15, ACCENT, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.5, y + 0.12, 2.6, 0.6, [[(cat, 12, ACCENT2, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 7.2, y + 0.12, 5.0, 0.6, [[(d, 12.5, GRAY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.0
footer(s, 2)
notes(s, (
    "Say: \u201CI picked these five because together they show range AND depth \u2014 agents, automation, tooling, "
    "live-data products, and gritty real-world reliability.\u201D\n"
    "\u2022 Don't dwell here. 15 seconds. It's a map, not a destination.\n"
    "\u2022 Optional: ask the room \u2014 \u201CTell me which one you want to go deep on first.\u201D It makes it a conversation, "
    "not a lecture, and shows confidence."
))

# ============================================================ 3. COMMANDCORE
deep_dive(3, "Standout \u00b7 01", "CommandCore",
    "github.com/rohith547/commandcore",
    "Imagine one chat window that understands what you want \u2014 and quietly dispatches the right specialist to do it.",
    "A control center where you type a request and it routes that intent to the right specialized AI agent \u2014 each with its own tools and memory. Telegram + inline keyboards, so there's zero learning curve.",
    "A router reads intent and hands off to an agent. State is persisted in a database, so it's restart-safe \u2014 reboot the server and it resumes exactly where it left off.",
    "Orchestration is the hard part: typed tools, persisted memory, and getting multiple agents to cooperate without stepping on each other. That's what separates a system from a demo.",
    ["Node/TS", "Prisma", "OpenAI", "Telegram", "Agents", "Docker"],
    ("HOOK (say it slowly, then pause):\n"
     "\u201CImagine one chat window that understands what you want \u2014 and quietly sends the right specialist to do it.\u201D\n\n"
     "STORY:\n"
     "\u2022 \u201CMost people bolt one chatbot onto one task. I built a control center instead.\u201D\n"
     "\u2022 \u201CYou type what you want. A router figures out your intent and hands it to a specialized agent \u2014 each "
     "one has its own tools and its own memory.\u201D\n"
     "\u2022 \u201CThe whole thing lives in Telegram with tap-button menus, so there's nothing to learn.\u201D\n\n"
     "THE HARD PART (slow down \u2014 this is where you sound like an expert):\n"
     "\u2022 \u201CThe real engineering isn't the chat \u2014 it's the orchestration. Typed tools so agents can't misfire. "
     "Persisted memory so context survives. And it's restart-safe \u2014 if the server reboots at 3am, it picks up "
     "exactly where it left off.\u201D\n\n"
     "LAND IT (then stop):\n"
     "\u201CThat's the difference between a demo you show once and a system you can actually depend on.\u201D\n\n"
     "IF THEY ASK:\n"
     "\u2022 Routing? \u2014 intent is classified, then dispatched to the matching agent handler.\n"
     "\u2022 Memory? \u2014 conversation + state persisted in the DB (Prisma), not held in RAM.\n"
     "\u2022 Why Telegram? \u2014 instant mobile reach, zero front-end to maintain, native button UX.")
)

# ============================================================ 4. BEACON
deep_dive(4, "Standout \u00b7 02", "Beacon",
    "github.com/rohith547/beacon",
    "Every morning a fresh AI news video gets made \u2014 and I barely touch it.",
    "An end-to-end content pipeline: it fetches the day's stories, verifies them, writes the script, generates a voice, renders the video, and publishes. The one human step is a single approval tap.",
    "Stages run in sequence \u2014 fetch \u2192 verify \u2192 script \u2192 voice \u2192 video \u2192 publish \u2014 each checkpointed so a failure never burns the whole run. Cost is capped per run.",
    "Trust. LLMs hallucinate, so I fact-check BEFORE generating a word, checkpoint every stage, and keep a human at the one decision that matters \u2014 the publish.",
    ["Python", "OpenAI", "TTS", "ffmpeg", "cron", "Human-in-loop"],
    ("HOOK:\n\u201CEvery morning, a fresh AI-news video gets made \u2014 and I barely touch it.\u201D\n\n"
     "STORY:\n"
     "\u2022 \u201CIt's a pipeline. It pulls the day's stories, verifies them before writing anything, scripts the "
     "segment, generates a voice, renders the video, and publishes it.\u201D\n"
     "\u2022 \u201CThe only thing I do is tap 'approve.' One human decision, in exactly the right place.\u201D\n\n"
     "THE HARD PART:\n"
     "\u2022 \u201CThe challenge with automation isn't speed \u2014 it's trust. Language models make things up. So I fact-"
     "check before I generate, not after. Every stage is checkpointed, so if step four fails, I don't lose steps "
     "one through three. And there's a hard cost cap per run.\u201D\n\n"
     "LAND IT:\n\u201CAutomation people trust isn't about removing the human \u2014 it's about putting the human in exactly "
     "the right place.\u201D\n\n"
     "IF THEY ASK:\n"
     "\u2022 How do you stop hallucinations? \u2014 verification gate before generation; unverified facts don't make it "
     "into the script.\n"
     "\u2022 What if a stage fails? \u2014 checkpointed + idempotent, so I re-run from the failed step, not the start.")
)

# ============================================================ 5. SKILLS HUB
deep_dive(5, "Standout \u00b7 03", "Skills Hub",
    "github.com/rohith547/skills-hub",
    "Here's how I build 75 apps without reinventing the wheel 75 times.",
    "25 reusable engineering 'skills' \u2014 my best patterns \u2014 distilled from across my projects, plus a browsable React app to explore them, with tests.",
    "Each skill is a documented folder wired directly into my AI command line, so every new build automatically follows my proven conventions instead of starting from scratch.",
    "It's compounding leverage: every project makes the next one faster. And it externalizes how I think, not just what I shipped \u2014 a system that builds systems.",
    ["React", "TypeScript", "Vite", "Vitest", "CLI", "Single-source"],
    ("HOOK:\n\u201CHere's how I build seventy-five apps without reinventing the wheel seventy-five times.\u201D\n\n"
     "STORY:\n"
     "\u2022 \u201CI took my best engineering patterns and turned them into 25 reusable skills.\u201D\n"
     "\u2022 \u201CThen I wired them straight into my AI command line \u2014 so every new project automatically follows my "
     "proven conventions from the first keystroke.\u201D\n"
     "\u2022 \u201COn top there's a clean app to browse them, and it's tested.\u201D\n\n"
     "WHY IT MATTERS (this is your differentiator \u2014 lean in):\n"
     "\u2022 \u201CThis is compounding leverage. Each project makes the next one faster. It's the reason I can ship at "
     "the volume I do without quality dropping.\u201D\n\n"
     "LAND IT:\n\u201CMost engineers ship code. I also ship the system that ships the code.\u201D\n\n"
     "IF THEY ASK:\n"
     "\u2022 What's a 'skill'? \u2014 a documented, reusable capability the AI loads automatically (e.g. build-app, "
     "deploy-app, test-coverage).\n"
     "\u2022 Single source of truth? \u2014 yes, one data file generates both the app and the installable skill folders.")
)

# ============================================================ 6. FINANCIAL INTELLIGENCE
deep_dive(6, "Standout \u00b7 04", "Financial Intelligence",
    "github.com/rohith547/smart-money-13f",
    "I turn raw market filings into signals a person can act on in seconds.",
    "A family of live dashboards: smart-money-13f tracks what big institutions buy, insider-tracker watches executives buying their own stock, macro-pulse scores the economy, fed-watch follows policy.",
    "Each pulls live financial data, normalizes messy filings, and renders a glanceable view \u2014 every component handles loading, error, empty and populated states.",
    "Financial data is messy and noisy. The value isn't fetching it \u2014 it's distilling it into something trustworthy and decision-ready that a human can read in seconds.",
    ["React", "TypeScript", "Tailwind", "Live data", "Dashboards", "Vercel"],
    ("HOOK:\n\u201CI turn raw market filings into signals a person can act on in seconds.\u201D\n\n"
     "STORY:\n"
     "\u2022 \u201CIt's a family of live dashboards. Smart-Money-13F shows what the big institutions are buying. Insider-"
     "Tracker flags when executives buy their own company's stock. Macro-Pulse rolls the whole economy into a "
     "handful of health scores. Fed-Watch tracks policy.\u201D\n"
     "\u2022 \u201CAll on live data \u2014 not static snapshots.\u201D\n\n"
     "THE HARD PART:\n"
     "\u2022 \u201CFinancial data is messy and noisy. The engineering is in normalizing filings and turning a wall of "
     "numbers into something you can glance at and trust. Every view handles loading, error, empty and "
     "populated \u2014 no half-built happy paths.\u201D\n\n"
     "LAND IT:\n\u201CData is everywhere. The value is in making it decision-ready.\u201D\n\n"
     "IF THEY ASK:\n"
     "\u2022 Where's the data from? \u2014 public filings / market feeds, normalized on ingest.\n"
     "\u2022 Why several apps? \u2014 each answers one question well; together they're a cohesive intelligence suite.")
)

# ============================================================ 7. RENT-TRACKER
deep_dive(7, "Standout \u00b7 05", "Rent-tracker",
    "private repo \u2014 live walkthrough on request",
    "I had a boring monthly question \u2014 who actually paid rent? \u2014 so I automated the whole answer.",
    "A Telegram bot logs into the bank, reads transactions, fuzzy-matches each deposit to the right tenant, and reports \u20187 of 7 tenants paid\u2019 every cycle. No spreadsheets.",
    "It scrapes the bank, matches deposits to tenants by name and amount, and pings me on Telegram \u2014 counting unique tenants, not raw transactions.",
    "The bank actively fights automation \u2014 sessions expire the instant you go headless. I solved it with fresh visible-browser logins, careful session handling, and idempotent runs so it never double-counts.",
    ["Python", "Playwright", "Telegram", "Fuzzy-match", "Idempotent", "Production"],
    ("HOOK:\n\u201CI had a boring monthly question \u2014 who actually paid rent? \u2014 so I automated the entire answer.\u201D\n\n"
     "STORY:\n"
     "\u2022 \u201CA bot logs into the bank, reads the transactions, matches each deposit to the right tenant, and every "
     "cycle just tells me: seven of seven tenants paid. No spreadsheet, no manual checking.\u201D\n\n"
     "THE HARD PART (this is your credibility moment \u2014 it sounds real because it IS):\n"
     "\u2022 \u201CThe bank fights automation. The moment you run it headless \u2014 invisibly \u2014 the session dies. I solved "
     "that with fresh visible-browser logins and proper session handling. And the runs are idempotent, so it "
     "never double-counts a payment.\u201D\n\n"
     "ON THE PRIVATE REPO (say this confidently):\n"
     "\u201CThis one's private because it touches real bank and tenant data \u2014 so instead of a link, I'll just walk "
     "you through the code live. Two minutes and you'll see exactly how it works.\u201D\n\n"
     "LAND IT:\n\u201CThis is the unglamorous, real-world engineering that actually saves time every month \u2014 and it's "
     "running in production right now.\u201D\n\n"
     "IF THEY ASK:\n"
     "\u2022 Why visible browser? \u2014 ICICI invalidates headless sessions instantly; a real browser session survives.\n"
     "\u2022 Matching accuracy? \u2014 fuzzy match on name + amount; counts unique tenants from the matched set.")
)

# ============================================================ 8. THE THROUGH-LINE
s = slide()
box(s, 0, 0, 0.18, SH.inches, fill=ACCENT)
text(s, 1.0, 1.55, 11.5, 1.4, [
    [("The through-line", 40, WHITE, True)],
    [("Five projects, one engineer, one consistent way of working.", 18, GRAY, False)],
], space_after=10)
pillars = [
    ("\U0001F3AF  Range + depth", "Agents, automation, tooling, live-data products, real-world reliability."),
    ("\u2699\uFE0F  Systems, not demos", "Restart-safe, idempotent, checkpointed, tested \u2014 built to run unattended."),
    ("\U0001F9E0  AI as a multiplier", "LLMs inside real products, with verification, cost caps and a human where it counts."),
]
x = 1.0
for t, d in pillars:
    box(s, x, 3.5, 3.7, 1.7, fill=PANEL, line=EDGE)
    text(s, x + 0.3, 3.68, 3.15, 1.4, [
        [(t, 15, ACCENT, True)],
        [(d, 12, GRAY, False)],
    ], space_after=5, line_spacing=1.06)
    x += 3.9
box(s, 1.0, 5.5, 11.3, 1.0, fill=PANEL, line=ACCENT, line_w=1.25)
text(s, 1.35, 5.66, 10.6, 0.8, [[
    ("So what?  ", 14, ACCENT, True),
    ("Point me at a problem and I'll turn it into a dependable, AI-powered product \u2014 fast, tested, and in "
     "production.", 14, WHITE, False),
]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s, 8)
notes(s, (
    "CLOSING \u2014 deliver this with conviction, slower than everything else:\n"
    "\u2022 \u201CIf you take one thing away: across all five of these, it's the same engineer working the same way.\u201D\n"
    "\u2022 \u201CRange across five domains. Depth in agents and automation. And the discipline to make it reliable \u2014 "
    "restart-safe, tested, running unattended.\u201D\n"
    "\u2022 \u201CI use AI as a multiplier, but the judgment \u2014 verifying before generating, capping cost, keeping a human "
    "where it matters \u2014 that's the engineering.\u201D\n\n"
    "FINAL LINE (look up, pause, then say it):\n"
    "\u201CPoint me at a problem, and I'll turn it into a dependable, AI-powered product \u2014 fast, tested, and in "
    "production. That's what I do.\u201D\n\n"
    "THEN: stop talking. Smile. Invite questions or offer the live walkthrough."
))

OUT = "Standout-Projects-Speaker-Deck.pptx"
prs.save(OUT)
print(f"Wrote {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
